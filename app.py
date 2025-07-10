import hashlib
from flask import Flask, request, jsonify, send_from_directory, render_template, redirect
from flask_mail import Mail, Message 
import requests
import os
from werkzeug.utils import secure_filename
from PyPDF2 import PdfReader
import json
from flask import session

# Importaciones opcionales para archivos especiales
try:
    from psd_tools import PSDImage
    HAS_PSD = True
except ImportError:
    HAS_PSD = False
    print("‚ö†Ô∏è  psd_tools no disponible. Archivos PSD no se procesar√°n.")

try:
    from pdf2image import convert_from_path
    HAS_PDF2IMAGE = True
except ImportError:
    HAS_PDF2IMAGE = False
    print("‚ö†Ô∏è  pdf2image no disponible. Archivos AI no se procesar√°n.")

try:
    from PIL import Image
    import pytesseract
    HAS_OCR = True
except ImportError:
    HAS_OCR = False
    print("‚ö†Ô∏è  PIL/pytesseract no disponible. OCR no funcionar√°.")

try:
    import docx
    HAS_DOCX = True
except ImportError:
    HAS_DOCX = False
    print("‚ö†Ô∏è  python-docx no disponible. Archivos DOCX no se procesar√°n.")

try:
    from pptx import Presentation
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False
    print("‚ö†Ô∏è  python-pptx no disponible. Archivos PPTX no se procesar√°n.")

try:
    import ezdxf
    HAS_EZDXF = True
except ImportError:
    HAS_EZDXF = False
    print("‚ö†Ô∏è  ezdxf no disponible. Archivos DWG/DXF no se procesar√°n.")

UPLOAD_FOLDER = 'uploads'
ALLOWED_EXTENSIONS = {
    'pdf', 'dwg', 'dxf', 'png', 'jpg', 'jpeg', 'bmp',
    'svg', 'tiff', 'webp', 'heic', 'gif', 'docx', 'pptx',
    'ai', 'psd', 'indd', 'txt'
}

app = Flask(__name__)
app.secret_key = 'una_clave_super_secreta'
app.config['MAIL_SERVER'] = 'smtp.gmail.com'
app.config['MAIL_PORT'] = 587
app.config['MAIL_USE_TLS'] = True
app.config['MAIL_USE_SSL'] = False
app.config['MAIL_USERNAME'] = 'tucorreo@gmail.com'
app.config['MAIL_PASSWORD'] = 'tu_contrase√±a_de_aplicacion'
app.config['MAIL_DEFAULT_SENDER'] = ('Tu App', 'tucorreo@gmail.com')

mail = Mail(app)
app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER
app.config['MAX_CONTENT_LENGTH'] = 16 * 1024 * 1024  # 16MB max file size

# Crear directorio de uploads
os.makedirs(UPLOAD_FOLDER, exist_ok=True)
USERS_FILE = os.path.join(UPLOAD_FOLDER, 'users.json')

def load_users():
    if not os.path.exists(USERS_FILE) or os.path.getsize(USERS_FILE) == 0:
        return []
    try:
        with open(USERS_FILE, 'r') as f:
            return json.load(f)
    except json.JSONDecodeError:
        return []

def save_users(users):
    with open(USERS_FILE, 'w') as f:
        json.dump(users, f, indent=4)

def allowed_file(filename):
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS

def buscar_serper(query, api_key):
    """Busca en Google usando Serper API"""
    url = "https://google.serper.dev/search"
    headers = {
        "X-API-KEY": api_key,
        "Content-Type": "application/json"
    }
    body = {"q": query}
    
    try:
        response = requests.post(url, headers=headers, json=body, timeout=10)
        print(f"üîç Serper API - Status: {response.status_code}")
        
        if response.status_code != 200:
            return None, f"Error HTTP: {response.status_code}"
        
        data = response.json()
        resultados = []
        organic = data.get("organic", [])
        
        if not organic:
            return None, "No se encontraron resultados."
        
        for r in organic[:5]:
            title = r.get("title", "")
            snippet = r.get("snippet", "")
            link = r.get("link", "")
            resultados.append(f"T√≠tulo: {title}\nResumen: {snippet}\nEnlace: {link}")
        
        return "\n\n".join(resultados), None
        
    except Exception as e:
        print(f"‚ùå Error en Serper: {e}")
        return None, f"Error conectando con Serper: {e}"

def llamar_a_ollama(prompt):
    """Llama al modelo Ollama local"""
    url = "http://127.0.0.1:11434/v1/chat/completions"
    system_message = {
        "role": "system",
        "content": "Eres un experto arquitecto y dise√±ador, pero tambi√©n sabes sobre todos los temas del mundo. "
                   "Ayuda con preguntas t√©cnicas y creativas. Habla como una persona normal, "
                   "responde en el idioma en que se te habla y no hagas saludos largos."
    }
    user_message = {
        "role": "user",
        "content": prompt
    }
    payload = {
        "model": "llama3",
        "messages": [system_message, user_message]
    }
    
    try:
        response = requests.post(url, json=payload, timeout=30)
        if response.status_code == 200:
            return response.json()['choices'][0]['message']['content']
        else:
            return f"Error del servidor del modelo: {response.status_code}"
    except Exception as e:
        return f"Error llamando al modelo: {e}"

def extract_text_from_file(filepath, filename):
    """
    Extrae texto de diferentes tipos de archivos
    """
    ext = filename.rsplit('.', 1)[1].lower() if '.' in filename else ''
    file_text = ""
    
    print(f"üìÑ Procesando archivo: {filename} (tipo: {ext})")
    
    try:
        if ext == 'txt':
            with open(filepath, 'r', encoding='utf-8') as f:
                file_text = f.read()
        
        elif ext == 'pdf':
            if not os.path.exists(filepath):
                return "[Error: Archivo PDF no encontrado]"
            
            reader = PdfReader(filepath)
            print(f"üìÑ PDF tiene {len(reader.pages)} p√°ginas")
            
            for i, page in enumerate(reader.pages[:5]):  # L√≠mite a 5 p√°ginas
                try:
                    page_text = page.extract_text()
                    if page_text:
                        file_text += f"--- P√°gina {i+1} ---\n{page_text}\n\n"
                except Exception as e:
                    print(f"Error en p√°gina {i+1}: {e}")
                    file_text += f"[Error leyendo p√°gina {i+1}]\n"

        elif ext in ('png', 'jpg', 'jpeg', 'bmp', 'gif', 'tiff', 'webp'):
            if not HAS_OCR:
                return "[OCR no disponible. Instala PIL y pytesseract]"
            
            try:
                img = Image.open(filepath)
                text = pytesseract.image_to_string(img, lang='spa+eng')
                file_text = text if text.strip() else "[No se detect√≥ texto en la imagen]"
            except Exception as e:
                file_text = f"[Error procesando imagen: {e}]"

        elif ext == 'docx':
            if not HAS_DOCX:
                return "[python-docx no disponible]"
            
            try:
                doc = docx.Document(filepath)
                paragraphs = []
                for para in doc.paragraphs:
                    if para.text.strip():
                        paragraphs.append(para.text)
                file_text = "\n".join(paragraphs)
            except Exception as e:
                file_text = f"[Error procesando DOCX: {e}]"

        elif ext == 'pptx':
            if not HAS_PPTX:
                return "[python-pptx no disponible]"
            
            try:
                ppt = Presentation(filepath)
                slides_text = []
                for i, slide in enumerate(ppt.slides):
                    slide_text = f"--- Diapositiva {i+1} ---\n"
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            slide_text += shape.text + "\n"
                    slides_text.append(slide_text)
                file_text = "\n".join(slides_text)
            except Exception as e:
                file_text = f"[Error procesando PPTX: {e}]"

        elif ext in ('dwg', 'dxf'):
            if not HAS_EZDXF:
                return "[ezdxf no disponible]"
            
            try:
                doc = ezdxf.readfile(filepath)
                msp = doc.modelspace()
                elements = []
                
                for entity in msp:
                    tipo = entity.dxftype()
                    if tipo == "LINE":
                        start = entity.dxf.start
                        end = entity.dxf.end
                        length = ((end[0] - start[0]) ** 2 + (end[1] - start[1]) ** 2) ** 0.5
                        elements.append(f"L√≠nea de {start} a {end}, longitud: {round(length, 2)}")
                    elif tipo == "TEXT":
                        elements.append(f"Texto: {entity.dxf.text}")
                    elif tipo == "MTEXT":
                        elements.append(f"Texto: {entity.text}")
                    elif tipo == "CIRCLE":
                        elements.append(f"C√≠rculo en {entity.dxf.center}, radio: {entity.dxf.radius}")
                
                file_text = "\n".join(elements) if elements else "[Archivo CAD sin elementos de texto]"
            except Exception as e:
                file_text = f"[Error procesando archivo CAD: {e}]"

        elif ext == 'psd':
            if not HAS_PSD:
                return "[psd_tools no disponible]"
            
            try:
                psd = PSDImage.open(filepath)
                img = psd.composite()
                if HAS_OCR:
                    text = pytesseract.image_to_string(img, lang='spa+eng')
                    file_text = text if text.strip() else "[No se detect√≥ texto en el PSD]"
                else:
                    file_text = "[PSD procesado pero OCR no disponible]"
            except Exception as e:
                file_text = f"[Error procesando PSD: {e}]"

        elif ext == 'ai':
            if not HAS_PDF2IMAGE:
                return "[pdf2image no disponible para archivos AI]"
            
            try:
                images = convert_from_path(filepath)
                if HAS_OCR:
                    ai_text = []
                    for i, img in enumerate(images):
                        text = pytesseract.image_to_string(img, lang='spa+eng')
                        if text.strip():
                            ai_text.append(f"--- P√°gina {i+1} ---\n{text}")
                    file_text = "\n".join(ai_text) if ai_text else "[No se detect√≥ texto en el archivo AI]"
                else:
                    file_text = "[Archivo AI procesado pero OCR no disponible]"
            except Exception as e:
                file_text = f"[Error procesando archivo AI: {e}]"

        elif ext == 'svg':
            try:
                with open(filepath, 'r', encoding='utf-8') as f:
                    svg_content = f.read()
                # Buscar texto en SVG
                import re
                text_matches = re.findall(r'<text[^>]*>(.*?)</text>', svg_content, re.DOTALL)
                if text_matches:
                    file_text = "\n".join(text_matches)
                else:
                    file_text = "[Archivo SVG sin texto]"
            except Exception as e:
                file_text = f"[Error procesando SVG: {e}]"

        elif ext == 'indd':
            file_text = "[Archivo InDesign - Exporta a PDF para procesarlo]"

        else:
            file_text = f"[Tipo de archivo '{ext}' no soportado]"

    except Exception as e:
        print(f"‚ùå Error general procesando {filename}: {e}")
        file_text = f"[Error procesando archivo: {e}]"

    print(f"üìÑ Texto extra√≠do: {len(file_text)} caracteres")
    return file_text

# --------------------------
# Rutas principales
# --------------------------

@app.route('/')
def home():
    return render_template('index.html', username=session.get('username'))

@app.route('/preguntar', methods=['POST'])
def preguntar():
    """Ruta para preguntas simples sin archivos"""
    data = request.json
    pregunta = data.get('pregunta')
    usar_serper = data.get('usar_serper', True)
    
    if not pregunta:
        return jsonify({"error": "No se recibi√≥ pregunta"}), 400

    api_key = "3afe5888cc608256a0ae579173d4fb0c7186a9d0"
    print(f"üîç Pregunta: {pregunta}")
    print(f"üîç Usar Serper: {usar_serper}")

    if usar_serper:
        print("üîç Buscando en Serper...")
        resultados, error = buscar_serper(pregunta, api_key)
        
        if resultados:
            print("‚úÖ Resultados encontrados en Serper")
            prompt = f"""Informaci√≥n actualizada de Internet:
{resultados}

Pregunta del usuario: {pregunta}

Responde usando esta informaci√≥n actualizada."""
        else:
            print(f"‚ùå Error en Serper: {error}")
            prompt = f"Pregunta: {pregunta}\nResponde con tu conocimiento."
    else:
        prompt = pregunta

    response_ollama = llamar_a_ollama(prompt)
    
    return jsonify({
        "respuesta": response_ollama,
        "busqueda_realizada": usar_serper
    })

@app.route('/api/generar-texto', methods=['POST'])
def generar_texto():
    """Ruta principal para texto + archivos"""
    
    # Obtener datos del request
    if request.is_json:
        data = request.get_json()
        prompt = data.get('prompt', '').strip()
        usar_serper = data.get('usar_serper', True)
        file = None
    else:
        prompt = request.form.get('prompt', '').strip()
        usar_serper = request.form.get('usar_serper', 'true').lower() == 'true'
        file = request.files.get('file')

    if not prompt:
        return jsonify({"error": "No se recibi√≥ prompt"}), 400

    print(f"üîç Prompt: {prompt}")
    print(f"üîç Usar Serper: {usar_serper}")
    print(f"üìÑ Archivo: {file.filename if file else 'No'}")

    # Buscar en Serper si est√° habilitado
    search_info = ""
    if usar_serper:
        api_key = "3afe5888cc608256a0ae579173d4fb0c7186a9d0"
        print("üîç Buscando en Serper...")
        resultados, error = buscar_serper(prompt, api_key)
        if resultados:
            print("‚úÖ Resultados encontrados en Serper")
            search_info = f"Informaci√≥n actualizada de Internet:\n{resultados}\n\n"
        else:
            print(f"‚ùå Error en Serper: {error}")

    # Procesar archivo si existe
    file_text = ""
    if file and file.filename:
        if allowed_file(file.filename):
            try:
                # Guardar archivo temporalmente
                filename = secure_filename(file.filename)
                filepath = os.path.join(app.config['UPLOAD_FOLDER'], filename)
                file.save(filepath)
                
                print(f"üìÑ Archivo guardado: {filepath}")
                
                # Extraer texto
                file_text = extract_text_from_file(filepath, filename)
                
                # Limpiar archivo temporal
                try:
                    os.remove(filepath)
                    print(f"üóëÔ∏è  Archivo temporal eliminado")
                except Exception as e:
                    print(f"‚ö†Ô∏è  Error eliminando archivo temporal: {e}")
                    
            except Exception as e:
                print(f"‚ùå Error procesando archivo: {e}")
                file_text = f"[Error procesando archivo: {e}]"
        else:
            file_text = f"[Tipo de archivo '{file.filename}' no permitido]"

    # Combinar informaci√≥n
    combined_info = []
    if search_info:
        combined_info.append(search_info)
    if file_text:
        combined_info.append(f"Contenido del archivo:\n{file_text}")
    
    if combined_info:
        combined_prompt = f"{''.join(combined_info)}\n\nPregunta del usuario: {prompt}"
    else:
        combined_prompt = prompt

    # Llamar a Ollama
    response_ollama = llamar_a_ollama(combined_prompt)
    
    return jsonify({
        'respuesta': response_ollama,
        'busqueda_realizada': usar_serper,
        'archivo_procesado': bool(file_text),
        'archivo_info': {
            'nombre': file.filename if file else None,
            'caracteres_extraidos': len(file_text) if file_text else 0
        }
    })

@app.route("/upload", methods=["POST"])
def upload_file():
    """Ruta espec√≠fica para subir archivos DXF y obtener info estructurada"""
    file = request.files.get("file")
    if not file or not file.filename:
        return jsonify({"error": "No se envi√≥ ning√∫n archivo."}), 400

    if not allowed_file(file.filename):
        return jsonify({"error": "Tipo de archivo no permitido."}), 400

    filename = secure_filename(file.filename)
    filepath = os.path.join(UPLOAD_FOLDER, filename)
    
    try:
        file.save(filepath)
        
        if filename.lower().endswith(".dxf"):
            if not HAS_EZDXF:
                return jsonify({"error": "ezdxf no disponible."}), 500
                
            doc = ezdxf.readfile(filepath)
            msp = doc.modelspace()
            entidades = []

            for entity in msp:
                tipo = entity.dxftype()
                if tipo == "LINE":
                    start = entity.dxf.start
                    end = entity.dxf.end
                    length = ((end[0] - start[0]) ** 2 + (end[1] - start[1]) ** 2) ** 0.5
                    entidades.append({
                        "tipo": "l√≠nea",
                        "inicio": [round(start[0], 2), round(start[1], 2)],
                        "fin": [round(end[0], 2), round(end[1], 2)],
                        "longitud": round(length, 2)
                    })
                elif tipo in ("TEXT", "MTEXT"):
                    texto = entity.dxf.text if tipo == "TEXT" else entity.text
                    entidades.append({
                        "tipo": "texto",
                        "contenido": texto
                    })

            return jsonify({"resultado": entidades})
        else:
            return jsonify({"mensaje": "Archivo recibido pero no es DXF."})

    except Exception as e:
        return jsonify({"error": f"No se pudo procesar el archivo: {str(e)}"}), 500

    finally:
        # Limpiar archivo temporal
        try:
            if os.path.exists(filepath):
                os.remove(filepath)
        except Exception as e:
            print(f"Error eliminando archivo temporal: {e}")

# --------------------------
# Rutas de autenticaci√≥n
# --------------------------

@app.route('/register', methods=['POST'])
def register():
    data = request.json
    email = data.get('email')
    password = data.get('password')
    name = data.get('name')

    if not name or not email or not password:
        return jsonify({"error": "Todos los campos son obligatorios."}), 400
    
    if len(password) < 6:
        return jsonify({"error": "La contrase√±a debe tener al menos 6 caracteres."}), 400

    users = load_users()

    if any(user['email'] == email for user in users):
        return jsonify({"error": "El correo ya est√° registrado."}), 409

    hashed_password = hashlib.sha256(password.encode('utf-8')).hexdigest()

    new_user = {
        "name": name,
        "email": email,
        "password": hashed_password
    }
    users.append(new_user)
    save_users(users)

    session['username'] = name

    # Enviar correo de bienvenida
    msg = Message(
        "Bienvenido a Nuestra P√°gina",
        recipients=[email]
    )
    msg.body = f"Hola {name}, gracias por registrarte!"
    msg.html = f"<h1>Hola {name}</h1><p>¬°Gracias por registrarte!</p>"

    try:
        mail.send(msg)
    except Exception as e:
        print(f"Error enviando correo: {e}")

    return jsonify({"message": "Usuario registrado correctamente."}), 201

@app.route('/login', methods=['POST'])
def login():
    if request.form:
        username = request.form.get('username')
        password = request.form.get('password')
    else:
        data = request.get_json()
        username = data.get('email')
        password = data.get('password')

    if not username or not password:
        return jsonify({"error": "Usuario/Email y contrase√±a son obligatorios."}), 400

    users = load_users()

    user_found = None
    for user in users:
        if user['email'] == username or user['name'] == username:
            user_found = user
            break

    if not user_found:
        return jsonify({"error": "Credenciales inv√°lidas."}), 401

    provided_password_hashed = hashlib.sha256(password.encode('utf-8')).hexdigest()

    if user_found['password'] == provided_password_hashed:
        session['username'] = user_found['name']

        if request.form:
            return redirect('/')
        else:
            return jsonify({
                "message": "Sesi√≥n iniciada.",
                "user": {
                    "name": user_found['name'],
                    "email": user_found['email']
                }
            }), 200
    else:
        return jsonify({"error": "Credenciales inv√°lidas."}), 401

@app.route('/logout')
def logout():
    session.pop('username', None)
    return redirect('/')

if __name__ == "__main__":
    app.run(debug=True, host='0.0.0.0', port=5050)